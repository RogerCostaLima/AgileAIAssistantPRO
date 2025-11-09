import pandas as pd
import io
from pptx import Presentation

# =====================
# EXPORTAÇÃO CSV/EXCEL
# =====================
def exportar_artefatos(resultados):
    """
    resultados = {
        "epic": "texto do epic",
        "feature": "texto da feature",
        "user_story": "texto user story",
        "task": "texto task"
    }
    """
    df = pd.DataFrame([resultados])
    return df

def baixar_excel(df, filename="artefatos.xlsx"):
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=False)
    buffer.seek(0)
    return buffer

# =====================
# PROCESSAR PPT PARA INSTRUÇÃO DA IA
# =====================
def extrair_texto_ppt(file):
    prs = Presentation(file)
    texto = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto += shape.text + "\n"
    return texto
